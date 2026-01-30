"""
Service clients for Azure resources.
"""
import os
import io
import json
import hashlib
import logging
from datetime import datetime, timedelta
from typing import Optional, Dict, Any, List
import pyodbc

from azure.servicebus import ServiceBusClient, ServiceBusMessage
from azure.cosmos import CosmosClient, PartitionKey
from azure.storage.blob import BlobServiceClient, generate_blob_sas, BlobSasPermissions, UserDelegationKey, ContentSettings
from azure.identity import DefaultAzureCredential

from pptx import Presentation
from pptx.util import Emu

logger = logging.getLogger(__name__)

# EMU (English Metric Units) per inch for PowerPoint measurements
EMU_PER_INCH = 914400


class ServiceBusService:
    """Service Bus client for job queue operations using managed identity."""

    def __init__(self):
        # Use managed identity - namespace should be fully qualified (e.g., "pptgen-dev-sb.servicebus.windows.net")
        self.namespace = os.environ.get("SERVICEBUS_NAMESPACE")
        self.queue_name = os.environ.get("SERVICEBUS_QUEUE_NAME", "ppt-generation-jobs")
        self._client = None

    def _get_client(self) -> ServiceBusClient:
        if self._client is None:
            if not self.namespace:
                raise ValueError("SERVICEBUS_NAMESPACE must be set for managed identity authentication")
            # Use managed identity via DefaultAzureCredential
            credential = DefaultAzureCredential()
            self._client = ServiceBusClient(
                fully_qualified_namespace=self.namespace,
                credential=credential
            )
        return self._client

    async def send_job(self, job_message: Dict[str, Any]) -> None:
        """Send a job to the processing queue."""
        client = self._get_client()
        with client.get_queue_sender(self.queue_name) as sender:
            message = ServiceBusMessage(
                body=json.dumps(job_message, default=str),
                message_id=job_message["jobId"],
                content_type="application/json",
                time_to_live=timedelta(hours=1)
            )
            sender.send_messages(message)
            logger.info(f"Job {job_message['jobId']} sent to queue")


class CosmosService:
    """Cosmos DB client for job state and caching."""

    def __init__(self):
        self.endpoint = os.environ.get("COSMOS_ENDPOINT")
        self.key = os.environ.get("COSMOS_KEY")
        self.database_name = os.environ.get("COSMOS_DATABASE", "ppt-generator")
        self._client = None
        self._database = None

    def _get_database(self):
        if self._database is None:
            if self.key:
                # Use key if provided
                self._client = CosmosClient(self.endpoint, self.key)
            else:
                # Use managed identity
                credential = DefaultAzureCredential()
                self._client = CosmosClient(self.endpoint, credential=credential)
            self._database = self._client.get_database_client(self.database_name)
        return self._database

    async def create_job(self, job_id: str, request_data: Dict[str, Any]) -> None:
        """Create a new job record."""
        db = self._get_database()
        container = db.get_container_client("jobs")

        job_doc = {
            "id": job_id,
            "partitionKey": job_id,  # Partition key for the container (matches /partitionKey path)
            "jobId": job_id,
            "status": "queued",
            "progress": 0,
            "currentStage": None,
            "requestId": request_data.get("requestId"),
            "callbackUrl": request_data.get("callbackUrl"),
            "createdAt": datetime.utcnow().isoformat(),
            "updatedAt": datetime.utcnow().isoformat(),
            "ttl": 604800  # 7 days
        }

        container.create_item(body=job_doc)
        logger.info(f"Job {job_id} created in Cosmos")

    async def get_job(self, job_id: str) -> Optional[Dict[str, Any]]:
        """Get job status by ID."""
        db = self._get_database()
        container = db.get_container_client("jobs")

        try:
            # Use cross-partition query to find job regardless of partition key issues
            query = "SELECT * FROM c WHERE c.id = @job_id"
            logger.info(f"Querying Cosmos for job {job_id}")
            items = list(container.query_items(
                query=query,
                parameters=[{"name": "@job_id", "value": job_id}],
                enable_cross_partition_query=True
            ))
            logger.info(f"Query returned {len(items)} items for job {job_id}")
            if items:
                return items[0]
            logger.warning(f"Job {job_id} not found in query - 0 results returned")
            return None
        except Exception as e:
            logger.exception(f"Job {job_id} query failed with exception: {e}")
            return None

    async def update_job(self, job_id: str, updates: Dict[str, Any]) -> None:
        """Update job status."""
        db = self._get_database()
        container = db.get_container_client("jobs")

        # First get the job using cross-partition query
        job = await self.get_job(job_id)
        if not job:
            logger.error(f"Cannot update job {job_id}: not found")
            return

        job.update(updates)
        job["updatedAt"] = datetime.utcnow().isoformat()

        # Use upsert to avoid partition key issues on replace
        container.upsert_item(body=job)

    async def check_cache(self, content_hash: str) -> Optional[Dict[str, Any]]:
        """Check if a similar request is cached."""
        db = self._get_database()
        container = db.get_container_client("cache")

        try:
            return container.read_item(item=content_hash, partition_key=content_hash)
        except Exception:
            return None

    async def cache_result(self, content_hash: str, result: Dict[str, Any]) -> None:
        """Cache a generation result."""
        db = self._get_database()
        container = db.get_container_client("cache")

        cache_doc = {
            "id": content_hash,
            "contentHash": content_hash,
            "outputUrl": result.get("downloadUrl"),
            "templateUsed": result.get("templateUsed"),
            "slideCount": result.get("slideCount"),
            "cachedAt": datetime.utcnow().isoformat(),
            "ttl": 86400  # 24 hours
        }

        container.upsert_item(body=cache_doc)

    async def log_error(self, job_id: str, error_data: Dict[str, Any]) -> None:
        """Log malformed data or errors to errors container."""
        db = self._get_database()
        container = db.get_container_client("errors")

        error_doc = {
            "id": f"{job_id}-{datetime.utcnow().timestamp()}",
            "jobId": job_id,
            "errorType": error_data.get("errorType"),
            "errorMessage": error_data.get("errorMessage"),
            "collectionTitle": error_data.get("collectionTitle"),
            "rawData": error_data.get("rawData"),
            "timestamp": datetime.utcnow().isoformat(),
            "ttl": 2592000  # 30 days
        }

        container.create_item(body=error_doc)

    @staticmethod
    def generate_content_hash(request_data: Dict[str, Any]) -> str:
        """Generate a hash for caching similar requests."""
        # Normalize the request for hashing (exclude requestId, callbackUrl)
        hashable = {
            "context": request_data.get("presentationContext", {}).get("context"),
            "dataCollections": [
                {"title": dc.get("title"), "data": dc.get("data")}
                for dc in request_data.get("dataCollections", [])
            ]
        }
        content = json.dumps(hashable, sort_keys=True, default=str)
        return hashlib.sha256(content.encode()).hexdigest()

    async def list_templates(self) -> list:
        """List all templates from Cosmos DB."""
        db = self._get_database()
        container = db.get_container_client("templates")

        try:
            query = "SELECT c.id, c.templateId, c.partitionKey, c.name FROM c"
            items = list(container.query_items(
                query=query,
                enable_cross_partition_query=True
            ))
            logger.info(f"Listed {len(items)} templates from Cosmos")
            return items
        except Exception as e:
            logger.exception(f"Failed to list templates from Cosmos: {e}")
            return []

    async def get_template(self, template_id: str) -> Optional[Dict[str, Any]]:
        """Get a single template from Cosmos DB by ID."""
        db = self._get_database()
        container = db.get_container_client("templates")

        try:
            query = "SELECT * FROM c WHERE c.id = @id OR c.templateId = @id"
            items = list(container.query_items(
                query=query,
                parameters=[{"name": "@id", "value": template_id}],
                enable_cross_partition_query=True
            ))

            if items:
                return items[0]
            return None
        except Exception as e:
            logger.exception(f"Failed to get template {template_id} from Cosmos: {e}")
            return None

    async def delete_template(self, template_id: str) -> None:
        """Delete a template from Cosmos DB."""
        db = self._get_database()
        container = db.get_container_client("templates")

        try:
            # First query to get the document with its partition key
            query = "SELECT * FROM c WHERE c.id = @id OR c.templateId = @id"
            items = list(container.query_items(
                query=query,
                parameters=[{"name": "@id", "value": template_id}],
                enable_cross_partition_query=True
            ))

            if not items:
                logger.warning(f"Template {template_id} not found in Cosmos for deletion")
                return

            for item in items:
                # Use the actual partition key from the document
                pk = item.get("partitionKey") or item.get("templateId") or item.get("id")
                doc_id = item.get("id")
                container.delete_item(item=doc_id, partition_key=pk)
                logger.info(f"Deleted template {template_id} (id={doc_id}, pk={pk}) from Cosmos")

        except Exception as e:
            logger.exception(f"Failed to delete template {template_id} from Cosmos: {e}")


class SQLTelemetryService:
    """Azure SQL client for telemetry and reporting."""

    def __init__(self):
        self.connection_string = os.environ.get("SQL_CONNECTION_STRING")

    def _get_connection(self):
        return pyodbc.connect(self.connection_string)

    async def log_request(self, job_id: str, request_data: Dict[str, Any]) -> None:
        """Log a new request to SQL for telemetry."""
        if not self.connection_string:
            logger.warning("SQL connection string not configured, skipping telemetry")
            return

        try:
            conn = self._get_connection()
            cursor = conn.cursor()

            requestor = request_data.get("requestorMetadata", {}) or {}
            context = request_data.get("presentationContext", {})

            cursor.execute("""
                INSERT INTO JobRequests (
                    JobId, RequestId, RequestorAppId, RequestorUserId, RequestorTenantId,
                    TemplatePreference, DataCollectionCount, ContextLength, Status
                ) VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?)
            """, (
                job_id,
                request_data.get("requestId"),
                requestor.get("appId"),
                requestor.get("userId"),
                requestor.get("tenantId"),
                context.get("templatePreference"),
                len(request_data.get("dataCollections", [])),
                len(context.get("context", "")),
                "queued"
            ))

            conn.commit()
            conn.close()
            logger.info(f"Telemetry logged for job {job_id}")
        except Exception as e:
            logger.error(f"Failed to log telemetry: {e}")

    async def update_job_status(
        self,
        job_id: str,
        status: str,
        template_used: Optional[str] = None,
        slide_count: Optional[int] = None,
        output_url: Optional[str] = None,
        error_code: Optional[str] = None,
        error_message: Optional[str] = None,
        duration_ms: Optional[int] = None
    ) -> None:
        """Update job status in SQL."""
        if not self.connection_string:
            return

        try:
            conn = self._get_connection()
            cursor = conn.cursor()

            cursor.execute("""
                UPDATE JobRequests SET
                    Status = ?,
                    TemplateUsed = COALESCE(?, TemplateUsed),
                    SlideCount = COALESCE(?, SlideCount),
                    OutputBlobUrl = COALESCE(?, OutputBlobUrl),
                    ErrorCode = COALESCE(?, ErrorCode),
                    ErrorMessage = COALESCE(?, ErrorMessage),
                    DurationMs = COALESCE(?, DurationMs),
                    CompletedAt = CASE WHEN ? IN ('completed', 'failed', 'completed_with_warnings')
                                       THEN GETUTCDATE() ELSE CompletedAt END
                WHERE JobId = ?
            """, (
                status, template_used, slide_count, output_url,
                error_code, error_message, duration_ms, status, job_id
            ))

            conn.commit()
            conn.close()
        except Exception as e:
            logger.error(f"Failed to update telemetry: {e}")


class BlobStorageService:
    """Azure Blob Storage client for templates and outputs."""

    def __init__(self):
        self.connection_string = os.environ.get("BLOB_CONNECTION_STRING")
        self.storage_account_name = os.environ.get("STORAGE_ACCOUNT_NAME")
        self.templates_container = os.environ.get("TEMPLATES_CONTAINER", "ppt-templates")
        self.output_container = os.environ.get("OUTPUT_CONTAINER", "ppt-outputs")
        self._client = None
        self._credential = None

    def _get_client(self) -> BlobServiceClient:
        if self._client is None:
            if self.connection_string:
                # Use connection string if available
                self._client = BlobServiceClient.from_connection_string(self.connection_string)
            elif self.storage_account_name:
                # Use managed identity with DefaultAzureCredential
                account_url = f"https://{self.storage_account_name}.blob.core.windows.net"
                self._credential = DefaultAzureCredential()
                self._client = BlobServiceClient(account_url, credential=self._credential)
            else:
                raise ValueError("Either BLOB_CONNECTION_STRING or STORAGE_ACCOUNT_NAME must be set")
        return self._client

    async def get_template_metadata(self, template_id: str) -> Optional[Dict[str, Any]]:
        """Get template metadata from blob storage."""
        client = self._get_client()
        container = client.get_container_client(self.templates_container)
        blob = container.get_blob_client(f"{template_id}/metadata.json")

        try:
            data = blob.download_blob().readall()
            return json.loads(data)
        except Exception as e:
            logger.warning(f"Template {template_id} metadata not found: {e}")
            return None

    async def list_active_templates(self) -> list:
        """List all active templates."""
        client = self._get_client()
        container = client.get_container_client(self.templates_container)

        templates = []
        for blob in container.list_blobs():
            if blob.name.endswith("metadata.json"):
                try:
                    blob_client = container.get_blob_client(blob.name)
                    data = json.loads(blob_client.download_blob().readall())
                    if data.get("active", True):
                        templates.append(data)
                except Exception as e:
                    logger.warning(f"Failed to load template metadata {blob.name}: {e}")

        return templates

    def generate_output_sas_url(self, job_id: str, expiry_hours: int = 24) -> str:
        """Generate a SAS URL for the output file."""
        client = self._get_client()
        blob_name = f"{job_id}/presentation.pptx"
        expiry = datetime.utcnow() + timedelta(hours=expiry_hours)

        if self.connection_string:
            # Use account key from connection string
            sas_token = generate_blob_sas(
                account_name=client.account_name,
                container_name=self.output_container,
                blob_name=blob_name,
                account_key=client.credential.account_key,
                permission=BlobSasPermissions(read=True),
                expiry=expiry
            )
        else:
            # Use user delegation key with managed identity
            start_time = datetime.utcnow() - timedelta(minutes=5)
            delegation_key = client.get_user_delegation_key(
                key_start_time=start_time,
                key_expiry_time=expiry
            )
            sas_token = generate_blob_sas(
                account_name=client.account_name,
                container_name=self.output_container,
                blob_name=blob_name,
                user_delegation_key=delegation_key,
                permission=BlobSasPermissions(read=True),
                expiry=expiry,
                start=start_time
            )

        return f"https://{client.account_name}.blob.core.windows.net/{self.output_container}/{blob_name}?{sas_token}"

    async def upload_template_metadata(self, template_id: str, metadata: Dict[str, Any]) -> None:
        """Upload generated metadata.json for a template."""
        client = self._get_client()
        container = client.get_container_client(self.templates_container)
        blob = container.get_blob_client(f"{template_id}/metadata.json")

        blob.upload_blob(
            json.dumps(metadata, indent=2),
            overwrite=True,
            content_settings=ContentSettings(content_type="application/json")
        )
        logger.info(f"Uploaded metadata for template {template_id}")

    def download_template(self, template_id: str) -> bytes:
        """Download a template file from blob storage."""
        client = self._get_client()
        container = client.get_container_client(self.templates_container)
        blob = container.get_blob_client(f"{template_id}/template.pptx")

        try:
            return blob.download_blob().readall()
        except Exception as e:
            # Try root-level template
            blob = container.get_blob_client(f"{template_id}.pptx")
            return blob.download_blob().readall()

    def upload_output(self, job_id: str, content: bytes) -> str:
        """Upload a generated presentation to the output container."""
        client = self._get_client()
        container = client.get_container_client(self.output_container)
        blob_name = f"{job_id}/presentation.pptx"
        blob = container.get_blob_client(blob_name)

        blob.upload_blob(
            content,
            overwrite=True,
            content_settings=ContentSettings(
                content_type="application/vnd.openxmlformats-officedocument.presentationml.presentation"
            )
        )
        logger.info(f"Uploaded output for job {job_id}")
        return blob_name


class TemplateIntrospectionService:
    """
    Introspects PowerPoint templates to extract layout and placeholder information.
    Generates metadata.json automatically when templates are uploaded.
    """

    # Mapping of PowerPoint placeholder types to our simplified types
    PLACEHOLDER_TYPE_MAP = {
        "CENTER_TITLE": "title",
        "TITLE": "title",
        "SUBTITLE": "subtitle",
        "BODY": "body",
        "OBJECT": "mixed",
        "PICTURE": "image",
        "CHART": "chart",
        "TABLE": "table",
        "FOOTER": "footer",
        "SLIDE_NUMBER": "slide_number",
        "DATE": "date",
    }

    # Content types each placeholder type can accept
    ACCEPTS_MAP = {
        "title": ["text"],
        "subtitle": ["text"],
        "body": ["text", "bullets"],
        "mixed": ["text", "bullets", "table", "chart", "image"],
        "image": ["image"],
        "chart": ["chart"],
        "table": ["table"],
        "footer": ["text"],
        "slide_number": ["text"],
        "date": ["text"],
    }

    def __init__(self, cosmos_service: Optional['CosmosService'] = None):
        self.cosmos = cosmos_service

    def introspect_template(self, template_bytes: bytes, template_id: str) -> Dict[str, Any]:
        """
        Analyze a PowerPoint template and extract its structure.

        Args:
            template_bytes: Raw bytes of the .pptx file
            template_id: Identifier for the template (derived from blob path)

        Returns:
            Complete metadata dictionary for the template
        """
        prs = Presentation(io.BytesIO(template_bytes))

        # Extract slide dimensions
        width_inches = round(prs.slide_width / EMU_PER_INCH, 2)
        height_inches = round(prs.slide_height / EMU_PER_INCH, 2)

        # Determine aspect ratio
        if abs(width_inches / height_inches - 16/9) < 0.1:
            aspect_ratio = "16:9"
        elif abs(width_inches / height_inches - 4/3) < 0.1:
            aspect_ratio = "4:3"
        else:
            aspect_ratio = f"{width_inches}:{height_inches}"

        # Extract layouts
        slide_layouts = []
        for idx, layout in enumerate(prs.slide_layouts):
            layout_info = self._introspect_layout(layout, idx)
            slide_layouts.append(layout_info)

        # Extract slide content and instructions from template slides
        slide_examples = []
        template_instructions = {
            "colorGuidelines": [],
            "formattingRules": [],
            "layoutGuidelines": [],
            "graphGuidelines": []
        }
        for idx, slide in enumerate(prs.slides):
            slide_info = self._extract_slide_content(slide, idx)
            if slide_info:
                slide_examples.append(slide_info)
                # Parse instructions from this slide
                self._parse_slide_instructions(slide_info, template_instructions)

        # Generate layout selection guide
        layout_selection_guide = self._generate_layout_selection_guide(slide_layouts)

        # Generate contentLayouts (filter to only layouts referenced in the guide)
        referenced_indices = {layout["layoutIndex"] for layout in layout_selection_guide.values() if "layoutIndex" in layout}
        content_layouts = [layout for layout in slide_layouts if layout["layoutIndex"] in referenced_indices]

        # Build metadata
        metadata = {
            "templateId": template_id,
            "version": "1.0.0",
            "name": self._generate_name(template_id),
            "description": f"Auto-generated metadata for template: {template_id}",
            "active": True,
            "useCases": [],
            "audienceTypes": ["general"],
            "bestFor": "General purpose presentations",
            "notRecommendedFor": [],
            "defaultBrandConfig": {
                "primaryColor": "#003366",
                "secondaryColor": "#E8F0F8",
                "baseColor": "#FFFFFF"
            },
            "dimensions": {
                "widthInches": width_inches,
                "heightInches": height_inches,
                "aspectRatio": aspect_ratio
            },
            "templateBlobPath": f"{template_id}/template.pptx",
            "slideLayouts": slide_layouts,
            "contentLayouts": content_layouts,
            "layoutSelectionGuide": layout_selection_guide,
            "slideExamples": slide_examples,
            "templateInstructions": template_instructions,
            "metadata": {
                "createdAt": datetime.utcnow().isoformat(),
                "updatedAt": datetime.utcnow().isoformat(),
                "createdBy": "template-introspection-service",
                "introspectedAt": datetime.utcnow().isoformat(),
                "tags": ["auto-generated"]
            }
        }

        return metadata

    def _introspect_layout(self, layout, layout_index: int) -> Dict[str, Any]:
        """Extract placeholder details from a slide layout."""
        placeholders = []

        for ph in layout.placeholders:
            ph_type_str = str(ph.placeholder_format.type).replace("PLACEHOLDER_TYPE.", "").split(" ")[0]
            our_type = self.PLACEHOLDER_TYPE_MAP.get(ph_type_str, "mixed")

            placeholder_info = {
                "idx": ph.placeholder_format.idx,
                "name": ph.name,
                "type": ph_type_str,
                "contentType": our_type,
                "accepts": self.ACCEPTS_MAP.get(our_type, ["text"]),
                "position": {
                    "leftInches": round(ph.left / EMU_PER_INCH, 2) if ph.left else 0,
                    "topInches": round(ph.top / EMU_PER_INCH, 2) if ph.top else 0,
                    "widthInches": round(ph.width / EMU_PER_INCH, 2) if ph.width else 0,
                    "heightInches": round(ph.height / EMU_PER_INCH, 2) if ph.height else 0,
                }
            }
            placeholders.append(placeholder_info)

        # Determine layout type based on placeholders present
        layout_type = self._infer_layout_type(placeholders, layout.name)

        # Determine if layout is duplicatable (content layouts usually are)
        duplicatable = layout_type in ["content", "chart", "table", "comparison", "mixed"]

        return {
            "layoutIndex": layout_index,
            "layoutName": layout.name,
            "layoutType": layout_type,
            "duplicatable": duplicatable,
            "placeholders": placeholders,
            "introspectedAt": datetime.utcnow().isoformat()
        }

    def _infer_layout_type(self, placeholders: List[Dict], layout_name: str) -> str:
        """Infer the layout type from placeholders and name."""
        name_lower = layout_name.lower()

        # Check for section dividers/headers first (various naming conventions)
        section_patterns = [
            "section header", "section break", "section divider", "section title",
            "sub title", "subtitle page", "chapter", "part title", "segment",
            "divider", "interstitial"
        ]
        # Check if any section pattern matches
        if any(pattern in name_lower for pattern in section_patterns):
            return "section_header"
        # Also check for "section" alone but not in "title and section"
        if "section" in name_lower and "and" not in name_lower:
            return "section_header"

        # Check name for title slides (main title, not section titles)
        if "title" in name_lower and "content" not in name_lower:
            # "Title Only" layouts are for charts, not section headers
            if "only" in name_lower:
                return "title_only"
            # Main title slide has minimal placeholders
            if len(placeholders) <= 2:
                return "title"

        if "blank" in name_lower:
            return "blank"
        if "comparison" in name_lower or "two content" in name_lower:
            return "comparison"
        if "chart" in name_lower:
            return "chart"
        if "table" in name_lower:
            return "table"

        # Check placeholder types
        content_types = [p["contentType"] for p in placeholders]

        if "chart" in content_types:
            return "chart"
        if "table" in content_types:
            return "table"
        if "image" in content_types and "body" not in content_types and "mixed" not in content_types:
            return "image"
        if "mixed" in content_types or "body" in content_types:
            return "content"
        if content_types.count("title") >= 1 and len(placeholders) <= 2:
            return "title"

        return "content"

    def _generate_layout_selection_guide(self, slide_layouts: List[Dict[str, Any]]) -> Dict[str, Any]:
        """
        Generate a layout selection guide that maps content types to appropriate layouts.

        Args:
            slide_layouts: List of introspected slide layouts

        Returns:
            Dictionary mapping content types to layout recommendations
        """
        guide = {}

        # Helper function to find best layout for a content type
        def find_layout(layout_type=None, name_contains=None, name_not_contains=None):
            for layout in slide_layouts:
                if layout_type and layout.get("layoutType") != layout_type:
                    continue
                layout_name_lower = layout.get("layoutName", "").lower()
                if name_contains:
                    name_patterns = name_contains if isinstance(name_contains, list) else [name_contains]
                    if not any(pattern in layout_name_lower for pattern in name_patterns):
                        continue
                if name_not_contains:
                    not_patterns = name_not_contains if isinstance(name_not_contains, list) else [name_not_contains]
                    if any(pattern in layout_name_lower for pattern in not_patterns):
                        continue
                return layout
            return None

        # Find title slide layout
        title_layout = find_layout(layout_type="title")
        if title_layout:
            guide["title_slide"] = {
                "layout": title_layout["layoutName"],
                "layoutIndex": title_layout["layoutIndex"],
                "when": "First slide of any presentation"
            }

        # Find section header layout
        section_layout = find_layout(layout_type="section_header")
        if section_layout:
            guide["section_header"] = {
                "layout": section_layout["layoutName"],
                "layoutIndex": section_layout["layoutIndex"],
                "when": "Separating major sections"
            }

        # Find bulleted text layout (content with body/mixed placeholders)
        content_layout = find_layout(layout_type="content", name_contains=["content", "bullet", "text"])
        if content_layout:
            guide["bulleted_text"] = {
                "layout": content_layout["layoutName"],
                "layoutIndex": content_layout["layoutIndex"],
                "when": "Bulleted lists, key points"
            }

        # Find chart/graph layout (title_only is versatile for full-width content)
        chart_layout = find_layout(layout_type="title_only")
        if not chart_layout:
            # Fallback to chart layout type
            chart_layout = find_layout(layout_type="chart")
        if chart_layout:
            guide["chart_or_graph"] = {
                "layout": chart_layout["layoutName"],
                "layoutIndex": chart_layout["layoutIndex"],
                "when": "Full-width charts, graphs, tables"
            }

        # Find text with graphic layout (two column, comparison, or layouts with "graphic" in name)
        graphic_layout = find_layout(name_contains=["two", "comparison", "graphic"])
        if graphic_layout:
            guide["text_with_graphic"] = {
                "layout": graphic_layout["layoutName"],
                "layoutIndex": graphic_layout["layoutIndex"],
                "when": "Text alongside chart or image"
            }

        # Find paragraph text layout (use title_only if available, it's versatile)
        paragraph_layout = find_layout(layout_type="title_only")
        if not paragraph_layout:
            # Fallback to blank or content layout
            paragraph_layout = find_layout(layout_type="blank")
        if not paragraph_layout:
            paragraph_layout = find_layout(layout_type="content")
        if paragraph_layout:
            guide["paragraph_text"] = {
                "layout": paragraph_layout["layoutName"],
                "layoutIndex": paragraph_layout["layoutIndex"],
                "when": "Paragraph text, flexible content"
            }

        return guide

    def _generate_name(self, template_id: str) -> str:
        """Generate a human-readable name from template ID."""
        return template_id.replace("-", " ").replace("_", " ").title()

    def _extract_slide_content(self, slide, slide_index: int) -> Optional[Dict[str, Any]]:
        """Extract text content from a template slide to identify instructions."""
        texts = []

        # Extract text from all shapes
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                texts.append(shape.text.strip())
            if hasattr(shape, "text_frame"):
                for paragraph in shape.text_frame.paragraphs:
                    para_text = "".join(run.text for run in paragraph.runs)
                    if para_text.strip():
                        texts.append(para_text.strip())

        if not texts:
            return None

        # Determine slide purpose based on content
        full_text = " ".join(texts).lower()
        slide_purpose = "content"

        if "read and delete" in full_text:
            slide_purpose = "instruction"
        elif "color" in full_text and ("palette" in full_text or "accent" in full_text):
            slide_purpose = "color_guide"
        elif "graph" in full_text or "chart" in full_text:
            slide_purpose = "chart_guide"
        elif "layout" in full_text and "select" in full_text:
            slide_purpose = "layout_guide"
        elif "icon" in full_text or "illustration" in full_text:
            slide_purpose = "asset_guide"

        # Get layout name if available
        layout_name = ""
        try:
            layout_name = slide.slide_layout.name
        except:
            pass

        return {
            "slideIndex": slide_index,
            "layoutName": layout_name,
            "purpose": slide_purpose,
            "textContent": texts[:20],  # Limit to first 20 text items
            "hasInstructions": "read and delete" in full_text or slide_purpose != "content"
        }

    def _parse_slide_instructions(self, slide_info: Dict[str, Any], instructions: Dict[str, List]) -> None:
        """Parse instructions from slide content and add to appropriate category."""
        if not slide_info.get("hasInstructions"):
            return

        full_text = " ".join(slide_info.get("textContent", [])).lower()
        purpose = slide_info.get("purpose", "")

        # Color guidelines
        if purpose == "color_guide" or ("accent" in full_text and "color" in full_text):
            color_rules = []
            if "20%" in full_text or "no more than" in full_text:
                color_rules.append("Accent colors should occupy no more than 20% of any layout")
            if "one" in full_text and "accent" in full_text:
                color_rules.append("Only one accent color should be used per slide/application")
            if "corporate blue" in full_text and "primary" in full_text:
                color_rules.append("Corporate Blue (#002856) is the primary corporate color and should be dominant")
            if color_rules:
                instructions["colorGuidelines"].extend(color_rules)

        # Graph/chart guidelines
        if purpose == "chart_guide" or ("graph" in full_text and "color" in full_text):
            graph_rules = []
            if "key" in full_text or "important" in full_text:
                graph_rules.append("Highlight key data with accent color")
            if "secondary" in full_text:
                graph_rules.append("Use primary brand color (Corporate Blue) for secondary data")
            if "tertiary" in full_text or "steel" in full_text:
                graph_rules.append("Use Steel gray for tertiary/less important data")
            if graph_rules:
                instructions["graphGuidelines"].extend(graph_rules)

        # Formatting rules
        if "text styling" in full_text or "specified sizes" in full_text:
            instructions["formattingRules"].append("Follow text styling and sizes consistently throughout presentation")

        if "divider" in full_text or purpose == "instruction":
            layout_name = slide_info.get("layoutName", "").lower()
            if "title" in full_text and "short" in full_text:
                if "divider" in layout_name or "sub" in layout_name:
                    instructions["formattingRules"].append("For section dividers: increase type size if copy is short")
                else:
                    instructions["formattingRules"].append("For title slides: increase type size if title is short")

        # Layout guidelines
        if purpose == "layout_guide":
            instructions["layoutGuidelines"].append("Multiple layout sets available: white/blue backgrounds with accent color options")

        # Deduplicate
        for key in instructions:
            instructions[key] = list(set(instructions[key]))

    async def save_to_cosmos(self, template_id: str, metadata: Dict[str, Any]) -> None:
        """Save template metadata to Cosmos DB for fast retrieval."""
        if not self.cosmos:
            logger.warning("Cosmos service not configured, skipping Cosmos save")
            return

        db = self.cosmos._get_database()
        container = db.get_container_client("templates")

        doc = {
            "id": template_id,
            "partitionKey": template_id,  # Required for Cosmos partition key
            "templateId": template_id,
            **metadata
        }

        try:
            container.upsert_item(body=doc)
            logger.info(f"Template metadata saved to Cosmos: {template_id}")
        except Exception as e:
            logger.error(f"Failed to save template to Cosmos: {e}")
